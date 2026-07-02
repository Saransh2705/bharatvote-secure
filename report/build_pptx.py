# -*- coding: utf-8 -*-
"""BharatVote presentation — Indian tricolor theme (navy/saffron/green),
matching the project report. For Saransh Kumar (BCA-V2, Chandigarh University)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY  = RGBColor(0x1A,0x2A,0x6C)
NAVY2 = RGBColor(0x22,0x30,0x7D)
SAFFR = RGBColor(0xFF,0x99,0x33)
GREEN = RGBColor(0x13,0x88,0x08)
GREEND= RGBColor(0x0E,0x6B,0x06)
INK   = RGBColor(0x0F,0x1E,0x3D)
MUTED = RGBColor(0x47,0x55,0x69)
TINT  = RGBColor(0xEE,0xF2,0xFF)
TINTS = RGBColor(0xFF,0xF7,0xED)
TINTG = RGBColor(0xF0,0xFD,0xF4)
WHITE = RGBColor(0xFF,0xFF,0xFF)
LIGHT = RGBColor(0xC7,0xD2,0xFE)
SAND  = RGBColor(0xFC,0xD9,0xB0)
LINE  = RGBColor(0xE2,0xE8,0xF0)

HEAD="Poppins"; BODY="Segoe UI"; SERIF="Georgia"; MONO="Consolas"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)
def no_line(s): s.line.fill.background()
def solid(s,c): s.fill.solid(); s.fill.fore_color.rgb=c
def grad(s,c1,c2,angle=25.0):
    try:
        s.fill.gradient(); st=s.fill.gradient_stops
        st[0].color.rgb=c1; st[0].position=0.0; st[1].color.rgb=c2; st[1].position=1.0
        try: s.fill.gradient_angle=angle
        except Exception: pass
    except Exception: solid(s,c1)
def rect(s,x,y,w,h,fill=None,gradient=None,rounded=False,radius=0.08,line=None,lw=1.0):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x),Inches(y),Inches(w),Inches(h))
    if rounded:
        try: shp.adjustments[0]=radius
        except Exception: pass
    if gradient: grad(shp,gradient[0],gradient[1],gradient[2] if len(gradient)>2 else 25.0)
    elif fill is not None: solid(shp,fill)
    else: shp.fill.background()
    if line is not None: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    else: no_line(shp)
    shp.shadow.inherit=False; return shp
def oval(s,x,y,w,h,fill=None,line=None,lw=1.5):
    shp=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is not None: solid(shp,fill)
    else: shp.fill.background()
    if line is not None: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    else: no_line(shp)
    shp.shadow.inherit=False; return shp
def arrow(s,x,y,w,h,fill):
    shp=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x),Inches(y),Inches(w),Inches(h))
    solid(shp,fill); no_line(shp); shp.shadow.inherit=False; return shp
def text(s,x,y,w,h,runs,size=16,color=INK,bold=False,font=BODY,align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP,sp_after=4,line_sp=1.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    paras=runs if isinstance(runs,list) else [runs]
    for i,para in enumerate(paras):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sp_after); p.space_before=Pt(0)
        try: p.line_spacing=line_sp
        except Exception: pass
        segs=para if isinstance(para,list) else [(para,color,bold,font,size)]
        for seg in segs:
            if isinstance(seg,tuple): txt,c,b,f,sz=(seg+(color,bold,font,size))[:5]
            else: txt,c,b,f,sz=seg,color,bold,font,size
            r=p.add_run(); r.text=txt
            r.font.size=Pt(sz); r.font.bold=b; r.font.name=f; r.font.color.rgb=c
    return tb
def banner(s,num,kicker,title):
    rect(s,0,0,13.333,1.18,gradient=(NAVY,NAVY2,20))
    oval(s,9.6,-1.0,3.2,3.2,fill=NAVY2)
    rect(s,0,1.18,13.333,0.07,fill=SAFFR)
    oval(s,0.55,0.26,0.66,0.66,fill=None,line=SAFFR,lw=2.0)
    text(s,0.55,0.26,0.66,0.66,str(num),size=22,color=WHITE,bold=True,font=HEAD,
         align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,1.45,0.20,11.2,0.3,kicker.upper(),size=10.5,color=SAND,bold=True,font=HEAD)
    text(s,1.42,0.46,11.4,0.6,title,size=25,color=WHITE,bold=True,font=HEAD)
def footer(s,n):
    rect(s,0.6,7.0,12.13,0.012,fill=LINE)
    text(s,0.6,7.08,9,0.35,[[("Saransh Kumar",MUTED,True,BODY,9),("  ·  023BCA110194  ·  Chandigarh University",MUTED,False,BODY,9)]],size=9)
    text(s,11.0,7.08,1.73,0.35,[[("Slide ",MUTED,False,BODY,9),(str(n),NAVY,True,BODY,9)]],size=9,align=PP_ALIGN.RIGHT)
def bullets(s,x,y,w,items,size=15,gap=8,marker=SAFFR):
    paras=[]
    for it in items:
        if isinstance(it,tuple): paras.append([("▸  ",marker,True,BODY,size),(it[0],INK,True,BODY,size),(it[1],INK,False,BODY,size)])
        else: paras.append([("▸  ",marker,True,BODY,size),(it,INK,False,BODY,size)])
    text(s,x,y,w,5,paras,size=size,sp_after=gap,line_sp=1.05)
def card(s,x,y,w,h,heading,body,hcolor=NAVY,fill=TINT,bar=None):
    rect(s,x,y,w,h,fill=fill,rounded=True,radius=0.06,line=LINE,lw=0.75)
    if bar: rect(s,x,y,0.08,h,fill=bar)
    text(s,x+0.22,y+0.16,w-0.44,0.4,heading,size=13.5,color=hcolor,bold=True,font=HEAD)
    text(s,x+0.22,y+0.62,w-0.44,h-0.7,body,size=11.5,color=MUTED,line_sp=1.02)
def style_cell(cell,txt,size=11,color=INK,bold=False,fill=WHITE,align=PP_ALIGN.LEFT,font=BODY):
    cell.fill.solid(); cell.fill.fore_color.rgb=fill
    cell.margin_left=Inches(0.12); cell.margin_right=Inches(0.08); cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
    cell.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=cell.text_frame.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=txt
    r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; r.font.name=font
def make_table(s,x,y,w,col_w,header,rows,hsize=11,bsize=10.5):
    nr=len(rows)+1; nc=len(header); gt=s.shapes.add_table(nr,nc,Inches(x),Inches(y),Inches(w),Inches(0.42*nr)).table
    gt.first_row=False; gt.horz_banding=False
    for j,cw in enumerate(col_w): gt.columns[j].width=Inches(cw)
    for j,h in enumerate(header): style_cell(gt.cell(0,j),h,size=hsize,color=WHITE,bold=True,fill=(NAVY if j%2==0 else NAVY2),font=HEAD)
    for i,row in enumerate(rows,start=1):
        f=WHITE if i%2 else TINT
        for j,v in enumerate(row): style_cell(gt.cell(i,j),v,size=bsize,color=(INK if j==0 else MUTED),bold=(j==0),fill=f)
    return gt

# 1 COVER
s=slide()
rect(s,0,0,13.333,7.5,gradient=(NAVY,NAVY2,35))
oval(s,9.3,-1.7,5.2,5.2,fill=SAFFR); oval(s,-1.8,4.6,4.4,4.4,fill=GREEN)
text(s,0.9,0.7,11.5,0.5,"CHANDIGARH UNIVERSITY",size=15,color=WHITE,bold=True,font=HEAD)
text(s,0.9,1.12,11.5,0.35,"University Institute of Computing  ·  Department of Computer Applications",size=10.5,color=LIGHT,font=BODY)
text(s,0.9,2.15,11,0.35,"FINAL-YEAR MAJOR PROJECT REPORT",size=12,color=SAND,bold=True,font=HEAD)
rect(s,0.95,2.62,2.4,0.09,fill=SAFFR)
text(s,0.88,2.8,11.5,1.2,"BharatVote",size=54,color=WHITE,bold=True,font=HEAD)
text(s,0.92,3.95,11.2,0.6,[[("A Secure Digital Voting Platform for the Indian Election System",WHITE,False,SERIF,18)]],size=18)
rect(s,0.9,4.9,6.9,2.0,fill=WHITE,rounded=True,radius=0.05)
rows=[("Submitted by","Saransh Kumar"),("University ID","023BCA110194"),("Programme","BCA-V2  ·  2023 – 2026"),("Mentor","Sahil Jindal — Quollabb")]
yy=5.08
for k,v in rows:
    text(s,1.2,yy,2.4,0.3,k,size=11.5,color=MUTED,font=BODY)
    text(s,3.6,yy,4.0,0.3,v,size=12.5,color=INK,bold=True,font=HEAD); yy+=0.45
text(s,0.9,7.0,11.5,0.35,"Industry mentorship by Quollabb, in partnership with Chandigarh University",size=10,color=LIGHT,font=SERIF)

# 2 INTRO
s=slide(); banner(s,1,"Introduction","Problem Statement"); footer(s,2)
text(s,0.6,1.55,12.1,1.0,"BharatVote is a secure, real-time digital voting platform for the Indian election system — unifying voter registration, verification, voting, and results in one application.",size=15.5,color=MUTED,font=SERIF,line_sp=1.1)
text(s,0.6,2.75,12,0.4,"A voting system must satisfy properties that are often in tension:",size=13.5,color=INK,bold=True,font=HEAD)
bullets(s,0.7,3.25,11.9,[("Integrity.  ","One vote per voter per election; recorded votes are immutable."),("Authentication.  ","Only verified, eligible citizens may vote."),("Verifiability.  ","A vote is private, yet the voter can confirm it was counted."),("Transparency.  ","Results are publicly auditable and timely.")],size=14,gap=9)
rect(s,0.6,6.05,12.13,0.8,fill=TINTS,rounded=True,radius=0.12); rect(s,0.6,6.05,0.1,0.8,fill=SAFFR)
text(s,0.95,6.18,11.6,0.6,[[("THE GOAL",SAFFR,True,HEAD,10)],[("Deliver these guarantees with a database-enforced security model and managed cloud services.",INK,False,BODY,13)]],size=13)

# 3 OBJECTIVES
s=slide(); banner(s,2,"Chapter One","Objectives"); footer(s,3)
objs=[("Domain model","Voters, elections, candidates, votes, booths, documents — a secure relational schema."),
("Registration","EPIC-based identity with an in-browser facial-capture step."),
("Secure voting","One-per-election integrity, encoded ballot, unique confirmation ID."),
("Verifiability","Independent vote verification by confirmation identifier."),
("Live results","Real-time vote counts and analytics."),
("DB-enforced security","Row Level Security + server-side privilege separation.")]
bars=[SAFFR,GREEN,NAVY,SAFFR,GREEN,NAVY]
for i,(h,b) in enumerate(objs):
    cx=0.6+(i%3)*4.11; cy=1.55+(i//3)*1.65; card(s,cx,cy,3.93,1.5,h,b,hcolor=NAVY,bar=bars[i])

# 4 SCOPE
s=slide(); banner(s,3,"Chapter One","Scope of the Project"); footer(s,4)
text(s,0.6,1.55,12.1,0.4,"In scope",size=15,color=GREEND,bold=True,font=HEAD)
bullets(s,0.7,2.0,11.9,["Full voter journey: register → verify → browse → vote → confirm → results","EPIC entry + in-browser facial capture; one-vote-per-election integrity","Live results with charts; vote verification; booth, documents, admin modules","Database-enforced privacy — sensitive data reachable only server-side"],size=14,gap=8,marker=GREEN)
text(s,0.6,4.55,12.1,0.4,"Out of scope (future work)",size=15,color=SAFFR,bold=True,font=HEAD)
bullets(s,0.7,5.0,11.9,["Live Aadhaar authentication & biometric face-matching (simulated here)","Formal end-to-end verifiable cryptographic voting protocol","SMS/OTP multi-factor and interactive booth maps"],size=14,gap=8,marker=SAFFR)

# 5 COMPARISON
s=slide(); banner(s,4,"Chapter Two","Existing Systems & Comparison"); footer(s,5)
make_table(s,0.6,1.65,12.13,[3.2,2.5,3.2,3.23],
    ["Property","Paper / EVM","Estonia i-Voting","BharatVote"],
    [["Remote participation","No","Yes","Yes (web)"],
     ["Identity verification","Manual / at booth","National eID","EPIC + facial capture"],
     ["One-vote integrity","Device / roll","Cryptographic","DB unique constraint"],
     ["Voter verifiability","VVPAT slip","Yes","Confirmation-ID lookup"],
     ["Real-time results","Counting day","After close","Live"],
     ["Unified voter services","No","Partial","Yes"]],hsize=11,bsize=10.5)
text(s,0.6,6.55,12.1,0.5,"BharatVote unifies registration, verification, voting, verification-of-vote, and results — enforcing integrity at the database layer.",size=12,color=MUTED,font=SERIF)

# 6 PROPOSED
s=slide(); banner(s,5,"Chapter Two","Proposed System — Overview"); footer(s,6)
text(s,0.6,1.55,12.1,0.9,"A three-layer, database-centric design: a React client, trusted Next.js server routes, and a policy-secured Supabase database.",size=14.5,color=MUTED,font=SERIF,line_sp=1.1)
tiers=[("Browser Client","Next.js · React · TS. Reads public data with a restricted anon key.",NAVY),
("Server Routes","Trusted API using the service-role key for all writes.",SAFFR),
("Supabase","PostgreSQL + Row Level Security + auto APIs + email.",GREEN)]
for i,(h,b,col) in enumerate(tiers):
    cx=0.6+i*4.12
    rect(s,cx,2.7,3.85,2.4,fill=TINT,rounded=True,radius=0.05,line=LINE,lw=0.75)
    rect(s,cx,2.7,3.85,0.62,fill=col,rounded=True,radius=0.05); rect(s,cx,3.0,3.85,0.32,fill=col)
    text(s,cx,2.78,3.85,0.5,h,size=15,color=WHITE,bold=True,font=HEAD,align=PP_ALIGN.CENTER)
    text(s,cx+0.25,3.5,3.35,1.5,b,size=12,color=MUTED,line_sp=1.1)
    if i<2: arrow(s,cx+3.88,3.75,0.2,0.3,SAFFR)
text(s,0.6,5.5,12.1,0.8,[[("Principle:  ",NAVY,True,HEAD,13)],[("The database is the source of truth; security is enforced where the data lives; the ballot log is append-only and one-per-voter.",MUTED,False,SERIF,13)]],size=13,line_sp=1.1)

# 7 TECH STACK
s=slide(); banner(s,6,"Implementation","Technology Stack"); footer(s,7)
stack=[("Next.js 15 + TypeScript","App Router hosts the React client and trusted server routes in one project — privilege separation without a separate backend.",NAVY),
("Supabase (PostgreSQL)","Managed database with Row Level Security, auto REST APIs, storage and keys — the whole backend as a service.",GREEN),
("Tailwind + shadcn/ui","Accessible, composable components and a government-portal design system, with Framer Motion transitions.",SAFFR),
("Recharts · react-webcam · Resend","Live result charts, in-browser facial capture, and transactional vote-confirmation email.",NAVY)]
for i,(h,b,col) in enumerate(stack):
    cx=0.6+(i%2)*6.17; cy=1.6+(i//2)*2.45; card(s,cx,cy,5.96,2.2,h,b,hcolor=col,bar=col)

# 8 ARCHITECTURE
s=slide(); banner(s,7,"Chapter Four","System Architecture"); footer(s,8)
def node(x,y,w,h,title,sub,col,tcol=WHITE):
    rect(s,x,y,w,h,fill=col,rounded=True,radius=0.08)
    paras=[[(title,tcol,True,HEAD,12.5)]]
    if sub: paras.append([(sub,(LIGHT if tcol==WHITE else MUTED),False,BODY,9.5)])
    text(s,x,y,w,h,paras,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,size=12.5)
node(0.7,2.5,2.8,1.5,"Browser Client","anon key · read only",NAVY)
node(4.1,2.5,3.2,1.5,"Next.js Server Routes","service-role · writes",SAFFR)
node(7.9,2.2,4.7,2.1,"Supabase","",GREEN)
text(s,7.9,2.35,4.7,0.4,"SUPABASE",size=11,color=WHITE,bold=True,font=HEAD,align=PP_ALIGN.CENTER)
rect(s,8.15,2.85,4.2,0.5,fill=WHITE,rounded=True,radius=0.1); text(s,8.15,2.9,4.2,0.4,"PostgreSQL · RLS · indexes",size=10.5,color=GREEND,bold=True,font=HEAD,align=PP_ALIGN.CENTER)
rect(s,8.15,3.45,2.0,0.45,fill=WHITE,rounded=True,radius=0.1); text(s,8.15,3.49,2.0,0.4,"Auto API",size=10,color=GREEND,align=PP_ALIGN.CENTER)
rect(s,10.35,3.45,2.0,0.45,fill=WHITE,rounded=True,radius=0.1); text(s,10.35,3.49,2.0,0.4,"Email",size=10,color=GREEND,align=PP_ALIGN.CENTER)
arrow(s,3.5,3.05,0.5,0.35,MUTED); arrow(s,7.3,3.05,0.55,0.35,MUTED)
text(s,0.7,4.15,2.8,0.4,"read public data",size=9.5,color=MUTED,align=PP_ALIGN.CENTER)
text(s,4.1,4.15,3.2,0.4,"register · vote · tally",size=9.5,color=MUTED,align=PP_ALIGN.CENTER)
text(s,0.6,5.35,12.1,1.0,"The browser holds only a restricted key and can read public electoral data but cannot write or read sensitive tables. All mutations flow through trusted server routes using the service-role key.",size=12.5,color=MUTED,font=SERIF,line_sp=1.1)

# 9 DATABASE
s=slide(); banner(s,8,"Chapter Four","Database Design"); footer(s,9)
tbls=[("users","id · voter_id · full_name · state · face_data",NAVY),
("elections","id · type · state · dates · status",SAFFR),
("candidates","id · election_id · name · party · votes",GREEN),
("votes","id · election_id · voter_id · candidate_id · confirmation_id",SAFFR),
("booths","id · booth_number · address · officer · lat/lng",NAVY),
("documents / announcements","official notices & publications",GREEN)]
for i,(h,b,col) in enumerate(tbls):
    cx=0.6+(i%3)*4.11; cy=1.65+(i//3)*1.75
    rect(s,cx,cy,3.9,1.5,fill=WHITE,rounded=True,radius=0.05,line=col,lw=1.25)
    rect(s,cx,cy,3.9,0.5,fill=col,rounded=True,radius=0.05); rect(s,cx,cy+0.22,3.9,0.28,fill=col)
    text(s,cx,cy+0.07,3.9,0.4,h,size=12.5,color=WHITE,bold=True,font=MONO,align=PP_ALIGN.CENTER)
    text(s,cx+0.2,cy+0.62,3.5,0.8,b,size=9.5,color=MUTED,font=MONO,line_sp=1.05)
text(s,0.6,5.4,12.1,0.9,"Seven normalized tables. The votes table is append-only and UNIQUE(election, voter) — double-voting is impossible by construction.",size=12.5,color=MUTED,font=SERIF)

# 10 FEATURES
s=slide(); banner(s,9,"Results","Key Features"); footer(s,10)
feats=["Voter registration","Facial capture","Identity verification","Browse elections","Secure voting","One-vote integrity","Vote encoding","Confirmation ID","Vote verification","Live results","Booth locator","Multi-level admin"]
for i,f in enumerate(feats):
    cx=0.6+(i%4)*3.05; cy=1.7+(i//4)*1.45
    rect(s,cx,cy,2.85,1.2,fill=TINT,rounded=True,radius=0.1,line=LINE,lw=0.75)
    oval(s,cx+0.22,cy+0.4,0.4,0.4,fill=[SAFFR,GREEN,NAVY,SAFFR][i%4])
    text(s,cx+0.8,cy,1.95,1.2,f,size=12.5,color=INK,bold=True,font=HEAD,anchor=MSO_ANCHOR.MIDDLE)

# 11 SECURITY
s=slide(); banner(s,10,"Chapter Four","Security — Row Level Security"); footer(s,11)
bullets(s,0.7,1.7,11.9,[("RLS on every table.  ","Access rules attached to the data itself, independent of client code."),("Public reads only.  ","Elections, candidates, booths, docs, announcements are readable by the anon key."),("Sensitive tables locked.  ","users and votes have NO anon policy — server-role access only.")],size=14.5,gap=11)
rect(s,0.6,3.85,12.13,1.15,fill=INK,rounded=True,radius=0.05)
text(s,0.9,4.02,11.5,0.9,[[("create policy \"read candidates\" on candidates for select using (true);",RGBColor(0x93,0xC5,0xFD),False,MONO,12)],[("-- users & votes: no anon policy → only the service-role key can touch them",RGBColor(0x6E,0xE7,0xB7),False,MONO,11.5)]],size=12,line_sp=1.35)
rect(s,0.6,5.35,12.13,1.05,fill=TINTG,rounded=True,radius=0.08); rect(s,0.6,5.35,0.1,1.05,fill=GREEN)
text(s,0.95,5.5,11.6,0.8,[[("GUARANTEE",GREEND,True,HEAD,11)],[("The browser can display public electoral data but can neither read voter identities/ballots nor write any table. Every mutation passes through trusted server code.",INK,False,BODY,12.5)]],size=12.5,line_sp=1.05)

# 12 VOTING FLOW
s=slide(); banner(s,11,"Chapter Four","Casting a Vote — Flow"); footer(s,12)
steps=[("Voter\nselects",NAVY),("/api/vote\nchecks",SAFFR),("append to\nvotes",GREEN),("tally +1\n+ email",NAVY),("confirmation\nID",SAFFR)]
x=0.7
for i,(label,col) in enumerate(steps):
    rect(s,x,2.6,1.95,1.3,fill=col,rounded=True,radius=0.1)
    text(s,x,2.6,1.95,1.3,label,size=12.5,color=WHITE,bold=True,font=HEAD,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    x+=1.95
    if i<len(steps)-1: arrow(s,x+0.05,3.05,0.42,0.4,MUTED); x+=0.52
text(s,0.6,4.4,12.1,1.2,"The server verifies the voter has not already voted in this election, appends the encoded ballot to the immutable votes table, increments the candidate's tally so results update live, sends a confirmation email, and returns a unique confirmation identifier for later verification.",size=13.5,color=MUTED,font=SERIF,line_sp=1.15)
rect(s,0.6,6.15,12.13,0.7,fill=TINTS,rounded=True,radius=0.1); rect(s,0.6,6.15,0.1,0.7,fill=SAFFR)
text(s,0.95,6.26,11.6,0.5,[[("UNIQUE(election, voter) makes double-voting impossible at the database level.",INK,False,BODY,12.5)]],size=12.5)

# 13 IMPLEMENTATION
s=slide(); banner(s,12,"Chapter Five","Implementation Highlights"); footer(s,13)
impl=[("Two-client data layer","Browser anon client for public reads; server service-role client for privileged writes.",NAVY),
("Trusted server routes","register, vote (cast + verify), elections, candidates — enforce all rules server-side.",SAFFR),
("Guided React flows","Registration (details → capture → done) and voting (select → confirm → done).",GREEN),
("One migration","Schema, RLS, indexes and seed data in a single version-controlled SQL file.",NAVY)]
for i,(h,b,col) in enumerate(impl):
    cx=0.6+(i%2)*6.17; cy=1.6+(i//2)*2.45; card(s,cx,cy,5.96,2.2,h,b,hcolor=col,bar=col)

# 14 TESTING
s=slide(); banner(s,13,"Chapter Six","Testing & Results"); footer(s,14)
make_table(s,0.6,1.65,12.13,[5.8,4.0,2.33],
    ["Test Case","Expected","Status"],
    [["Register voter","Record created + session","Pass"],
     ["Cast vote","Recorded + confirmation ID","Pass"],
     ["Live tally","Candidate count +1","Pass"],
     ["Double-vote blocked","Second attempt rejected","Pass"],
     ["Verify by confirmation ID","Vote confirmed","Pass"],
     ["RLS: anon read users/votes","Access denied","Pass"],
     ["Production build","17 routes, types valid","Pass"]],hsize=11.5,bsize=11)
text(s,0.6,6.55,12.1,0.5,"All test cases passed; verified end-to-end on the live cloud database, and the production build compiled cleanly.",size=12.5,color=MUTED,font=SERIF)

# 15 CONCLUSION
s=slide(); banner(s,14,"Chapter Nine","Conclusion & Future Scope"); footer(s,15)
text(s,0.6,1.55,12.1,1.0,"A secure, verifiable, real-time digital voting platform — with integrity and privacy enforced at the database layer — was built and verified end-to-end on a live cloud backend, with no separately operated server.",size=14.5,color=INK,font=SERIF,line_sp=1.15)
text(s,0.6,3.0,12,0.4,"Future scope",size=15,color=SAFFR,bold=True,font=HEAD)
fs=[("Aadhaar & biometrics",NAVY),("E2E verifiable crypto",GREEN),("SMS / OTP MFA",SAFFR),("Tamper-evident audit",NAVY),("Interactive booth maps",GREEN),("Multi-language & a11y",SAFFR)]
for i,(f,col) in enumerate(fs):
    cx=0.6+(i%3)*4.11; cy=3.55+(i//3)*1.15
    rect(s,cx,cy,3.93,0.95,fill=TINT,rounded=True,radius=0.1,line=LINE,lw=0.75); rect(s,cx,cy,0.09,0.95,fill=col)
    text(s,cx+0.3,cy,3.5,0.95,f,size=12.5,color=INK,bold=True,font=HEAD,anchor=MSO_ANCHOR.MIDDLE)

# 16 REFERENCES
s=slide(); banner(s,15,"Bibliography","References"); footer(s,16)
refs=["Next.js — Documentation. nextjs.org/docs","Supabase — Documentation. supabase.com/docs","PostgreSQL — Row Security Policies. postgresql.org/docs","React — react.dev  ·  Tailwind CSS — tailwindcss.com","shadcn/ui — ui.shadcn.com  ·  Recharts — recharts.org","Resend — Email API. resend.com/docs","Election Commission of India — EVM & VVPAT. eci.gov.in","Republic of Estonia — i-Voting. e-estonia.com","Springall et al., \"Security Analysis of the Estonian Internet Voting System,\" ACM CCS 2014","Specter et al., \"The Ballot is Busted Before the Blockchain: A Security Analysis of Voatz,\" USENIX Security 2020"]
paras=[[("[%d]  "%(i+1),SAFFR,True,BODY,12),(r,INK,False,BODY,12)] for i,r in enumerate(refs)]
text(s,0.7,1.65,12.0,5.2,paras,size=12,sp_after=6,line_sp=1.05)

# 17 THANK YOU
s=slide()
rect(s,0,0,13.333,7.5,gradient=(NAVY,NAVY2,35)); oval(s,9.6,-1.9,5.6,5.6,fill=SAFFR); oval(s,-1.6,4.4,4.6,4.6,fill=GREEN)
rect(s,5.4,2.55,2.5,0.1,fill=SAFFR)
text(s,0.9,2.85,11.5,1.2,"Thank You",size=58,color=WHITE,bold=True,font=HEAD,align=PP_ALIGN.CENTER)
text(s,0.9,4.2,11.5,0.5,"Questions & Discussion",size=18,color=LIGHT,font=SERIF,align=PP_ALIGN.CENTER)
text(s,0.9,5.5,11.5,0.8,[[("Saransh Kumar",WHITE,True,HEAD,15)],[("UID 023BCA110194  ·  BCA-V2  ·  Chandigarh University  ·  Mentor: Sahil Jindal (Quollabb)",LIGHT,False,BODY,12)]],size=15,align=PP_ALIGN.CENTER)

import os
target=r"D:\MegaP\bharatvote-secure\report\BharatVote-Presentation.pptx"
try:
    prs.save(target); print("SAVED:",target,"|",len(prs.slides._sldIdLst),"slides")
except PermissionError:
    alt=r"D:\MegaP\bharatvote-secure\report\BharatVote-Presentation-new.pptx"
    prs.save(alt); print("LOCKED — SAVED:",alt,"|",len(prs.slides._sldIdLst),"slides")
